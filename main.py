
# import Presentation class
# from pptx library
from pptx import Presentation
from pathlib import Path
from os.path import isfile, join
import os
from alive_progress import alive_bar
from tkinter import Tk
from tkinter.filedialog import askdirectory

# Method to get content from slides
def getPowerpointContent(path):
    ppt = Presentation(path)
    content = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text
                # print(shape.text)
                #
    return content


# creating an object
# powerpoint_directory = "C:/Users/mitow/Desktop/source/GRACE HYMNS/"
powerpoint_directory = askdirectory(title="Select Folder source of PPT")
global_path = Path(powerpoint_directory)
file_type = input("What file type are you converting (pptx or ppt)?")
file_list = [str(pp) for pp in global_path.glob("*." + file_type)]

print(type(file_list))
corpus = [str(f) for f in os.listdir(powerpoint_directory) if not f.startswith('.') and isfile(join(powerpoint_directory, f))]

with alive_bar(len(corpus)) as bar:
    for filename in corpus:
        print(filename)
        path = powerpoint_directory + "/" + filename
        print(path)
        file_content = getPowerpointContent(path)
        f = open(powerpoint_directory + "/output/" + filename.split(".")[0] + ".txt", "w+", encoding="utf-8")
        f.write(str(file_content))
        f.close()
        bar()

print("Done")
